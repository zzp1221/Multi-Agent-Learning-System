"""用于早期集成的结构化沙箱资产生成。"""

from __future__ import annotations

import json
import base64
import logging
import re
from dataclasses import asdict
from pathlib import Path
from typing import Any

from pydantic import BaseModel, ConfigDict, Field

from src.ai_modules.config import get_sandbox_root, get_settings
from src.ai_modules.generation.content_chain import (
    ContentGenerationChain,
    GeneratedMindMap,
    GeneratedSectionBundle,
)
from src.ai_modules.generation.html_ppt_builder import HtmlPptDeckBuilder
from src.ai_modules.models import (
    VideoGenerationTaskPayload,
    VideoSandboxArtifact,
)
from src.ai_modules.runtime import SystemSnapshot

LOGGER = logging.getLogger(__name__)
DEFAULT_VIDEO_TTS_MAX_CHARS = 260


class GeneratedAsset(BaseModel):
    """写入沙箱存储的已生成资产的元数据。"""

    asset_type: str = Field(alias="assetType")
    title: str
    summary: str
    display_mode: str = Field(alias="displayMode")
    file_name: str = Field(default="", alias="fileName")
    local_path: str | None = Field(default=None, alias="localPath")
    preview_text: str = Field(alias="previewText")
    mime_type: str | None = Field(default=None, alias="mimeType")
    inline_content: str | None = Field(default=None, alias="inlineContent")
    language: str | None = None
    explanation: str | None = None
    thumbnail_path: str | None = Field(default=None, alias="thumbnailPath")
    thumbnail_file_name: str | None = Field(default=None, alias="thumbnailFileName")
    thumbnail_mime_type: str | None = Field(default=None, alias="thumbnailMimeType")
    duration_seconds: int | None = Field(default=None, alias="durationSeconds")
    video_style: str | None = Field(default=None, alias="videoStyle")
    knowledge_point: str | None = Field(default=None, alias="knowledgePoint")
    generated_by: str | None = Field(default=None, alias="generatedBy")
    content_origin: str | None = Field(default=None, alias="contentOrigin")
    provider: str | None = None
    model: str | None = None
    agent_name: str | None = Field(default=None, alias="agentName")
    evidence_ids: list[str] = Field(default_factory=list, alias="evidenceIds")
    fallback: bool | None = None
    from_cache: bool = Field(default=False, alias="fromCache")

    model_config = ConfigDict(populate_by_name=True)


class SectionPlan(BaseModel):
    """已生成教学资产中的规划章节。"""

    title: str
    objective: str
    source_titles: list[str] = Field(default_factory=list, alias="sourceTitles")

    model_config = ConfigDict(populate_by_name=True)


class ResourceGenerationService:
    """将结构化资产输出写入本地沙箱目录。"""

    def __init__(
        self,
        sandbox_root: Path | None = None,
        content_chain: ContentGenerationChain | None = None,
    ) -> None:
        self.sandbox_root = sandbox_root or get_sandbox_root()
        self.content_chain = content_chain or ContentGenerationChain()

    @staticmethod
    def _video_tts_text(script_text: str, max_chars: int) -> str:
        normalized = " ".join(script_text.split())
        limit = max(200, max_chars or DEFAULT_VIDEO_TTS_MAX_CHARS)
        if len(normalized) <= limit:
            return normalized
        head = normalized[:limit]
        sentence_end = max(head.rfind("。"), head.rfind("！"), head.rfind("？"), head.rfind("."))
        if sentence_end >= int(limit * 0.6):
            return head[: sentence_end + 1].strip()
        return head.rstrip("，,；;、 ") + "。"

    async def build_asset(
        self,
        *,
        asset_type: str,
        params: dict,
        snapshot: SystemSnapshot,
    ) -> GeneratedAsset:
        builder_map = {
            "DOCUMENT": self._build_document,
            "READING": self._build_reading,
            "SLIDES": self._build_slides,
            "MINDMAP": self._build_mindmap,
            "CODE": self._build_code,
            "VIDEO": self.build_video_asset,
        }
        builder = builder_map.get(asset_type)
        if builder is None:
            raise ValueError(f"Unsupported assetType: {asset_type}")
        return await builder(params=params, snapshot=snapshot)

    async def build_video_asset(
        self,
        *,
        params: dict,
        snapshot: SystemSnapshot,
    ) -> GeneratedAsset:
        topic = self._display_topic(params)
        task_id = str(params.get("taskId") or "video-task")
        style = str(params.get("style") or "hybrid")
        duration_target = self._normalize_duration_seconds(params.get("duration") or params.get("durationTarget"))
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", []) if isinstance(retrieval, dict) else []
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        script_payload = await self.content_chain.generate_video_script_async(
            title=f"{topic}教学视频",
            topic=topic,
            snapshot=generation_snapshot,
            sources=sources,
            duration_seconds=duration_target,
            style=style,
        )
        task_dir = self.sandbox_root / f"video_{self._safe_task_id(task_id)}"
        task_dir.mkdir(parents=True, exist_ok=True)

        script_json_path = task_dir / "script.json"
        script_text_path = task_dir / "script.txt"
        audio_path = task_dir / "speech.mp3"
        thumbnail_path = task_dir / "thumbnail.svg"

        script_json_path.write_text(
            json.dumps(script_payload.model_dump(by_alias=True), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        script_text_path.write_text(script_payload.full_text, encoding="utf-8")

        tts_audio_bytes = params.get("tts_audio_bytes")
        if not tts_audio_bytes or not isinstance(tts_audio_bytes, bytes) or len(tts_audio_bytes) < 100:
            from src.ai_modules.llms.mimo_client import MiMoClient

            try:
                settings = get_settings()
                mimo_client = MiMoClient(timeout_seconds=settings.tts_timeout_seconds)
                tts_text = self._video_tts_text(script_payload.full_text, settings.video_tts_max_chars)
                LOGGER.info(
                    "Generating video TTS audio: task_id=%s text_chars=%s timeout_seconds=%s",
                    task_id,
                    len(tts_text),
                    settings.tts_timeout_seconds,
                )
                tts_audio_bytes = await mimo_client.synthesize_speech(
                    text=tts_text,
                    style_description="用清晰自然的语速播报，声音沉稳专业，适合教学场景",
                    voice="mimo_default",
                    audio_format="mp3",
                )
            except Exception as exc:
                raise RuntimeError(
                    f"Video TTS generation failed: {type(exc).__name__}: {exc}"
                ) from exc
            params["tts_audio_bytes"] = tts_audio_bytes
        audio_path.write_bytes(tts_audio_bytes)

        thumbnail_path.write_text(
            self._build_video_thumbnail_svg(
                title=script_payload.title,
                topic=topic,
                style=style,
            ),
            encoding="utf-8",
        )

        artifact = VideoSandboxArtifact(
            taskDir=str(task_dir),
            scriptJsonPath=str(script_json_path),
            scriptTextPath=str(script_text_path),
            audioPath=str(audio_path),
            finalVideoPath=None,
            thumbnailPath=str(thumbnail_path),
            durationSeconds=script_payload.total_duration,
            videoStyle=style,
            previewText=script_payload.full_text[:100],
            summaryText=f"{topic} 教学视频脚本与语音已生成，等待浏览器本地渲染。",
            audioBase64=base64.b64encode(tts_audio_bytes).decode("utf-8"),
            audioFormat="mp3",
            avatarDataUrl="/dh_live/assets/combined_data.json.gz",
        )
        settings = get_settings()
        params["videoSandboxArtifact"] = artifact.model_dump(by_alias=True)
        params["videoGenerationTask"] = VideoGenerationTaskPayload(
            status="completed",
            title=script_payload.title,
            topic=topic,
            script=script_payload,
            durationSeconds=artifact.duration_seconds,
            videoStyle=style,
            ttsProvider=settings.tts_provider,
            avatarProvider="browser_dh_live_mini",
            generationParams={
                "durationTarget": script_payload.total_duration,
                "style": style,
            },
        ).model_dump(by_alias=True)
        return GeneratedAsset(
            assetType="VIDEO",
            title=script_payload.title,
            summary=artifact.summary_text,
            displayMode="VIDEO_PLAYER",
            fileName="browser-rendered.webm",
            localPath=None,
            previewText=artifact.preview_text,
            mimeType="video/webm",
            thumbnailPath=str(thumbnail_path),
            thumbnailFileName="thumbnail.svg",
            thumbnailMimeType="image/svg+xml",
            durationSeconds=artifact.duration_seconds,
            videoStyle=style,
            knowledgePoint=topic,
        )

    async def _build_document(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        display_topic = self._display_topic(params)
        title = f"{display_topic}导学文档"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        section_plans = self._plan_document_sections(
            params=params,
            snapshot=snapshot,
            sources=sources,
        )
        generated_sections = await self.content_chain.generate_document_sections(
            title=title,
            topic=display_topic,
            snapshot=generation_snapshot,
            section_plans=[plan.model_dump(by_alias=True) for plan in section_plans],
            sources=sources,
        )
        content = self._render_document_markdown(
            title=title,
            params=params,
            snapshot=snapshot,
            section_plans=section_plans,
            sources=sources,
            generated_sections=generated_sections,
        )
        file_name = self._scoped_file_name("document_guide", "md", params)
        path = self._write_text(file_name, content)
        return GeneratedAsset(
            assetType="DOCUMENT",
            title=title,
            summary="基于检索证据生成的结构化课程导学文档",
            displayMode="MARKDOWN_CARD",
            fileName=file_name,
            localPath=str(path),
            previewText=self._build_preview_text(section_plans),
            mimeType="text/markdown; charset=UTF-8",
            inlineContent=content,
        )

    def _plan_document_sections(
        self,
        *,
        params: dict,
        snapshot: SystemSnapshot,
        sources: list[dict[str, Any]],
    ) -> list[SectionPlan]:
        topic = self._display_topic(params)
        source_titles = [str(item.get("title", "未知来源")) for item in sources[:5]]
        learner_gap = ", ".join(snapshot.knowledge_gaps[:2]) or "暂无明确薄弱点"
        return [
            SectionPlan(
                title="一、核心概念与学习目标",
                objective=f"帮助学生围绕 `{topic}` 建立基础概念框架。",
                sourceTitles=source_titles[:2],
            ),
            SectionPlan(
                title="二、关键原理与判断方法",
                objective=f"从 `{snapshot.current_course}` 视角解释原理与判断依据。",
                sourceTitles=source_titles[:3],
            ),
            SectionPlan(
                title="三、典型误区与辨析",
                objective=f"聚焦薄弱点 `{learner_gap}`，说明常见误区与纠偏方式。",
                sourceTitles=source_titles[1:4],
            ),
            SectionPlan(
                title="四、练习建议与复习路径",
                objective="给出可执行的练习顺序、复习策略和下一步建议。",
                sourceTitles=source_titles[:2],
            ),
        ]

    def _render_document_markdown(
        self,
        *,
        title: str,
        params: dict,
        snapshot: SystemSnapshot,
        section_plans: list[SectionPlan],
        sources: list[dict[str, Any]],
        generated_sections: GeneratedSectionBundle,
    ) -> str:
        topic = self._display_topic(params)
        lines = [
            f"# {title}",
            "",
            "## 文档概览",
            f"本文围绕 `{topic}` 组织内容，采用“概念 -> 原理 -> 误区 -> 练习”的生成链路展开。",
            "",
            "## 生成大纲",
            *[f"- {plan.title}: {plan.objective}" for plan in section_plans],
        ]

        for section_index, (plan, generated_section) in enumerate(
            zip(section_plans, generated_sections.sections, strict=False),
            start=1,
        ):
            lines.extend(
                [
                    "",
                    f"## {generated_section.title or plan.title}",
                    generated_section.body,
                    "",
                    "### 学习提示",
                    *generated_section.tips,
                ]
            )

        return "\n".join(lines)

    def render_section_paragraph(
        self,
        *,
        plan: SectionPlan | dict[str, Any],
        snapshot: SystemSnapshot | dict[str, Any],
        topic: str,
        section_index: int,
    ) -> str:
        snapshot_student_level = self._snapshot_value(snapshot, "student_level")
        snapshot_gaps = self._snapshot_list(snapshot, "knowledge_gaps")
        paragraph_by_section = {
            1: (
                f"`{topic}` 是当前知识点中最先需要建立的概念锚点。"
                f" 对于 `{snapshot_student_level}` 水平的学生，建议先回答“它是什么、解决什么问题、和相邻概念有什么区别”。"
            ),
            2: (
                f"理解 `{topic}` 时，不要只记结论，更要抓住判断条件与使用边界。"
                " 学习时可以结合课程中的典型例题，把原理和题目条件一一对应起来。"
            ),
            3: (
                f"从当前画像看，学生在 `{', '.join(snapshot_gaps) or '暂无明确薄弱点'}` 上更容易出错。"
                f" 因此本节重点解释 `{topic}` 与易混概念之间的边界，以及常见错因。"
            ),
            4: (
                f"完成 `{topic}` 的学习后，建议立即安排小规模练习，并把错题回流到薄弱点记录中。"
                " 先做基础题验证概念，再做综合题训练迁移。"
            ),
        }
        objective = plan.get("objective", "") if isinstance(plan, dict) else plan.objective
        return paragraph_by_section.get(section_index, objective)

    def render_section_tips(
        self,
        *,
        plan: SectionPlan | dict[str, Any],
        snapshot: SystemSnapshot | dict[str, Any],
        section_index: int,
    ) -> list[str]:
        current_chapter = self._snapshot_value(snapshot, "current_chapter")
        tips_by_section = {
            1: [f"- 先用一句话复述概念，再和 `{current_chapter}` 中相邻知识点做区分。"],
            2: ["- 尝试把判断条件写成 2-3 条清单，减少“只会背不会用”的情况。"],
            3: ["- 做题时先圈出限制条件，再判断是否满足使用前提。"],
            4: ["- 练习后及时记录错因，并把错因映射回当前画像薄弱点。"],
        }
        objective = plan.get("objective", "") if isinstance(plan, dict) else plan.objective
        return tips_by_section.get(section_index, [f"- {objective}"])

    def render_section_citations(self, *, plan: SectionPlan | dict[str, Any]) -> list[str]:
        source_titles = (
            plan.get("sourceTitles", [])
            if isinstance(plan, dict)
            else plan.source_titles
        )
        if not source_titles:
            return ["- [来源] 当前为回退内容，待补真实检索来源。"]
        return [
            f"- [来源{index}] {source_title}"
            for index, source_title in enumerate(source_titles, start=1)
        ]

    def _build_preview_text(self, section_plans: list[SectionPlan]) -> str:
        if not section_plans:
            return "已生成课程资源"
        return f"已生成结构化文档，共 {len(section_plans)} 个章节"

    def _snapshot_value(
        self,
        snapshot: SystemSnapshot | dict[str, Any],
        key: str,
    ) -> Any:
        if isinstance(snapshot, dict):
            return snapshot.get(key, "")
        return getattr(snapshot, key)

    def _snapshot_list(
        self,
        snapshot: SystemSnapshot | dict[str, Any],
        key: str,
    ) -> list[str]:
        value = self._snapshot_value(snapshot, key)
        return list(value) if isinstance(value, list) else []

    def _build_generation_snapshot(
        self,
        *,
        params: dict[str, Any],
        snapshot: SystemSnapshot,
    ) -> dict[str, Any]:
        profile = params.get("profile", {}) if isinstance(params.get("profile", {}), dict) else {}
        base = asdict(snapshot)
        base["preferred_resource_types"] = list(profile.get("preferredResourceTypes", []))
        base["learning_goal"] = (
            profile.get("learningGoal")
            or (profile.get("currentGoal", {}) or {}).get("shortTerm")
            or ""
        )
        return base

    async def _build_reading(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        display_topic = self._display_topic(params)
        title = f"{display_topic}延伸阅读"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        generated_reading = await self.content_chain.generate_reading_asset(
            title=title,
            topic=display_topic,
            snapshot=generation_snapshot,
            sources=sources,
        )
        content = "\n".join([f"# {generated_reading.title}", "", generated_reading.body])
        file_name = self._scoped_file_name("reading_material", "md", params)
        path = self._write_text(file_name, content)
        return GeneratedAsset(
            assetType="READING",
            title=generated_reading.title,
            summary=generated_reading.summary,
            displayMode="MARKDOWN_CARD",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_reading.title,
            mimeType="text/markdown; charset=UTF-8",
            inlineContent=content,
        )

    async def _build_slides(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        display_topic = self._display_topic(params)
        title = f"{display_topic}PPT大纲"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        topic = display_topic
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        if self._requires_slide_outline_confirmation(params):
            outline_content = await self._generate_slide_outline_markdown(
                title=title,
                topic=topic,
                params=params,
                snapshot=generation_snapshot,
                sources=sources,
            )
            return GeneratedAsset(
                assetType="SLIDES",
                title=title,
                summary="PPT 大纲已生成，等待用户确认后再生成演示文件",
                displayMode="SLIDE_OUTLINE_CONFIRMATION",
                fileName="",
                localPath=None,
                previewText=title,
                mimeType="text/markdown; charset=UTF-8",
                inlineContent=outline_content,
            )

        deck_html, slide_count = self._generate_html_ppt_with_omni(
            title=title, topic=topic, snapshot=snapshot, sources=sources, params=params
        )
        file_name = self._scoped_file_name("slides", "html", params)
        path = self._write_text(file_name, deck_html)
        return GeneratedAsset(
            assetType="SLIDES",
            title=title,
            summary=f"html-ppt 生成的可演示 HTML 课件 ({slide_count} 页)",
            displayMode="DOWNLOAD_CARD",
            fileName=file_name,
            localPath=str(path),
            previewText=f"HTML PPT 课件 · {slide_count} 页 · {topic}",
            mimeType="text/html; charset=UTF-8",
        )

    def _requires_slide_outline_confirmation(self, params: dict[str, Any]) -> bool:
        return not self._has_confirmed_slide_outline(params)

    def _has_confirmed_slide_outline(self, params: dict[str, Any]) -> bool:
        confirmed_outline = self._confirmed_slide_outline_text(params)
        return self._truthy(self._confirmed_slide_outline_value(params)) and bool(confirmed_outline)

    @staticmethod
    def _confirmed_slide_outline_value(params: dict[str, Any] | None) -> Any:
        if not isinstance(params, dict):
            return None
        if params.get("confirmedSlideOutline") is not None:
            return params.get("confirmedSlideOutline")
        learning_context = params.get("learningContext")
        if isinstance(learning_context, dict):
            return learning_context.get("confirmedSlideOutline")
        return None

    @staticmethod
    def _confirmed_slide_outline_text(params: dict[str, Any] | None) -> str:
        if not isinstance(params, dict):
            return ""
        top_level = str(params.get("confirmedSlideOutlineText") or "").strip()
        if top_level:
            return top_level
        learning_context = params.get("learningContext")
        if isinstance(learning_context, dict):
            return str(learning_context.get("confirmedSlideOutlineText") or "").strip()
        return ""

    async def _generate_slide_outline_markdown(
        self,
        *,
        title: str,
        topic: str,
        params: dict[str, Any],
        snapshot: dict[str, Any],
        sources: list[dict[str, Any]],
    ) -> str:
        generated_slides = await self.content_chain.generate_slides_asset(
            title=title,
            topic=topic,
            snapshot=snapshot,
            sources=sources,
        )
        lines = [f"# {generated_slides.title}", "", generated_slides.summary, ""]
        for index, slide in enumerate(generated_slides.slides, start=1):
            lines.extend(
                [
                    f"## {index}. {slide.title}",
                    *[f"- {bullet}" for bullet in slide.bullets],
                    "",
                    f"讲解备注：{slide.speaker_notes}",
                    "",
                ]
            )
        confirmed_outline = self._confirmed_slide_outline_text(params)
        if confirmed_outline:
            lines.extend(["## 用户确认的大纲", confirmed_outline])
        return "\n".join(lines).strip() + "\n"

    @staticmethod
    def _truthy(value: Any) -> bool:
        return str(value).strip().lower() in {"1", "true", "yes", "y", "confirmed"}

    def _generate_html_ppt_with_omni(
        self,
        *,
        title: str,
        topic: str,
        snapshot: SystemSnapshot,
        sources: list[dict[str, Any]],
        params: dict[str, Any] | None = None,
    ) -> tuple[str, int]:
        slides = self._generate_validated_slides_with_omni(
            title=title,
            topic=topic,
            snapshot=snapshot,
            sources=sources,
            params=params,
        )
        deck_html = HtmlPptDeckBuilder().render(
            title=title,
            topic=topic,
            course=str(snapshot.current_course),
            slides=slides,
        )
        return deck_html, len(slides) + 3

    def _generate_pptx_with_omni(
        self,
        *,
        title: str,
        topic: str,
        snapshot: SystemSnapshot,
        sources: list[dict[str, Any]],
        params: dict[str, Any] | None = None,
    ) -> bytes:
        """通过 MiMo-V2-Omni 生成 PPTX；失败时显式中止。"""
        slides = self._generate_validated_slides_with_omni(
            title=title,
            topic=topic,
            snapshot=snapshot,
            sources=sources,
            params=params,
        )
        pptx_bytes = self._build_pptx_bytes(
            title=title,
            topic=topic,
            slides=slides,
            course=str(snapshot.current_course),
        )
        if not pptx_bytes:
            raise RuntimeError("PPTX generation failed: python-pptx produced no file")
        return pptx_bytes

    def _generate_validated_slides_with_omni(
        self,
        *,
        title: str,
        topic: str,
        snapshot: SystemSnapshot,
        sources: list[dict[str, Any]],
        params: dict[str, Any] | None = None,
    ) -> list[dict[str, Any]]:
        """通过 MiMo-V2-Omni 生成并校验结构化幻灯片内容。"""
        settings = get_settings()
        if not settings.mimo_api_key:
            raise RuntimeError("PPTX generation failed: missing MIMO_API_KEY")

        try:
            from src.ai_modules.llms.mimo_client import MiMoClient

            source_texts = "\n".join(
                f"- {s.get('title', 'unknown')}: {s.get('evidence', '')[:200]}"
                for s in sources[:4]
            )
            confirmed_outline = self._confirmed_slide_outline_text(params)
            outline_instruction = (
                f"\n用户已确认以下 PPT 大纲，请优先按这个结构生成：\n{confirmed_outline}\n"
                if confirmed_outline
                else ""
            )
            prompt = (
                f"请为教学主题「{topic}」生成一份完整的 PPT 内容，用于 {snapshot.current_course} 课程。\n"
                f"学生水平: {snapshot.student_level}，学习风格: {snapshot.preferred_style}。\n"
                f"参考来源:\n{source_texts}\n\n"
                f"{outline_instruction}"
                "这些课程、学生画像和来源信息只用于生成判断，最终幻灯片正文不要展示课程、学生水平、学习风格、参考来源或证据说明等元信息。\n"
                "请以JSON格式输出，包含以下字段：\n"
                '{{"slides":[{{"slideTitle":"标题","bullets":["要点1","要点2"],"speakerNotes":"讲解备注"}}]}}\n'
                "要求：6-10页幻灯片，每页3-5个要点，speakerNotes用中文写50-100字的讲解说明。"
                "首尾页分别为标题页和总结页。仅输出JSON。"
            )

            client = MiMoClient()
            response = client.omni_chat_sync(
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=4096,
            )
            slide_data = client.extract_json(response)
            if not slide_data:
                raise RuntimeError("PPTX generation failed: MiMo returned no parseable slide JSON")

            slides = self._validate_slide_deck(slide_data.get("slides", []))
            return slides
        except RuntimeError:
            raise
        except Exception as exc:
            raise RuntimeError(f"PPTX generation failed: {type(exc).__name__}: {exc}") from exc

    @staticmethod
    def _validate_slide_deck(slides: Any) -> list[dict[str, Any]]:
        if not isinstance(slides, list) or not slides:
            raise RuntimeError("PPTX generation failed: MiMo returned an empty slide list")
        if len(slides) < 6 or len(slides) > 10:
            raise RuntimeError("PPTX generation failed: MiMo slide deck must contain 6-10 content slides")

        normalized_slides: list[dict[str, Any]] = []
        for index, slide_info in enumerate(slides, start=1):
            if not isinstance(slide_info, dict):
                raise RuntimeError(f"PPTX generation failed: slide {index} is not a structured object")
            slide_title = str(slide_info.get("slideTitle") or slide_info.get("title") or "").strip()
            raw_bullets = slide_info.get("bullets", [])
            bullets = [str(item).strip() for item in raw_bullets if str(item).strip()] if isinstance(raw_bullets, list) else []
            speaker_notes = str(slide_info.get("speakerNotes") or slide_info.get("speaker_notes") or "").strip()
            if not slide_title:
                raise RuntimeError(f"PPTX generation failed: slide {index} is missing a title")
            if not speaker_notes:
                raise RuntimeError(f"PPTX generation failed: slide {index} is missing speaker notes")
            bullets = ResourceGenerationService._normalize_slide_bullets(
                slide_title=slide_title,
                bullets=bullets,
                speaker_notes=speaker_notes,
            )
            normalized_slides.append(
                {
                    "slideTitle": slide_title,
                    "bullets": bullets,
                    "speakerNotes": speaker_notes,
                }
            )
        return normalized_slides

    @staticmethod
    def _normalize_slide_bullets(*, slide_title: str, bullets: list[str], speaker_notes: str) -> list[str]:
        normalized = bullets[:5]
        candidates = [slide_title]
        candidates.extend(
            item.strip()
            for item in re.split(r"[。！？!?；;\n\r]+", speaker_notes)
            if item.strip()
        )
        candidates.append(speaker_notes)

        for candidate in candidates:
            if len(normalized) >= 3:
                break
            bullet = ResourceGenerationService._compact_slide_bullet(candidate)
            if bullet and bullet not in normalized:
                normalized.append(bullet)

        while len(normalized) < 3:
            normalized.append(f"{ResourceGenerationService._compact_slide_bullet(slide_title)} {len(normalized) + 1}")
        return normalized

    @staticmethod
    def _compact_slide_bullet(text: str) -> str:
        compacted = re.sub(r"\s+", " ", text).strip(" -:：，,。.;；")
        if len(compacted) <= 40:
            return compacted
        return compacted[:40].rstrip()

    @staticmethod
    def _build_pptx_bytes(
        *,
        title: str,
        topic: str,
        slides: list[dict[str, Any]],
        course: str,
    ) -> bytes | None:
        """使用 python-pptx 在内存中构建 PPTX 文件。"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return None

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── 标题页 ──
        title_slide_layout = prs.slide_layouts[0]  # 标题页布局
        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        subtitle_placeholder.text = f"{course}\n{topic}"

        # ── 内容页 ──
        for slide_info in slides:
            slide_title = slide_info.get("slideTitle", slide_info.get("title", ""))
            bullets = slide_info.get("bullets", [])
            speaker_notes_text = slide_info.get("speakerNotes", slide_info.get("speaker_notes", ""))

            bullet_layout = prs.slide_layouts[1]  # 标题 + 内容
            slide = prs.slides.add_slide(bullet_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            # 添加要点
            body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body_shape and bullets:
                tf = body_shape.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0
                    p.font.size = Pt(24)

            # 讲解备注
            if speaker_notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = str(speaker_notes_text)

        # ── 总结页 ──
        summary_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "总结与回顾"
        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            summary_points = [
                f"主题: {topic}",
                f"共 {len(slides)} 个内容页",
                "请结合课堂讨论加深理解",
            ]
            for i, point in enumerate(summary_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = point
                p.font.size = Pt(24)

        import io
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    @staticmethod
    def _count_pptx_slides(pptx_bytes: bytes) -> int:
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(pptx_bytes))
            return len(prs.slides)
        except Exception:
            return 0

    def _write_bytes(self, file_name: str, data: bytes) -> Path:
        self.sandbox_root.mkdir(parents=True, exist_ok=True)
        path = self.sandbox_root / file_name
        path.write_bytes(data)
        return path

    def _display_topic(self, params: dict) -> str:
        learning_context = params.get("learningContext", {})
        nested_active_step = None
        if isinstance(learning_context, dict):
            nested_active_step = learning_context.get("activeLearningStep")

        strict_candidates = [
            params.get("explicitUserTopic"),
            learning_context.get("explicitUserTopic") if isinstance(learning_context, dict) else None,
            params.get("activeLearningStepTitle"),
            learning_context.get("activeLearningStepTitle") if isinstance(learning_context, dict) else None,
            nested_active_step,
            params.get("topic"),
            params.get("keyPoints"),
            params.get("knowledgePoint"),
            learning_context.get("knowledgePoint") if isinstance(learning_context, dict) else None,
            learning_context.get("chapter") if isinstance(learning_context, dict) else None,
            learning_context.get("course") if isinstance(learning_context, dict) else None,
        ]
        for candidate in strict_candidates:
            value = self._normalize_topic_candidate(candidate)
            if self._is_real_topic(value):
                return value
        for candidate in (params.get("rewrittenQuery"), params.get("query")):
            value = self._normalize_topic_candidate(candidate)
            if self._is_real_topic(value) and not self._looks_like_resource_command(value):
                return value
        raise RuntimeError("缺少资源生成真实主题：请提供当前学习阶段或明确的资源主题")

    @staticmethod
    def _normalize_topic_candidate(value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, dict):
            for key in ("title", "name", "label"):
                nested = str(value.get(key) or "").strip()
                if nested:
                    return nested
            return ""
        if isinstance(value, (list, tuple, set)):
            parts = [str(item).strip() for item in value if str(item).strip()]
            return "、".join(parts[:6])
        return str(value).strip()

    @staticmethod
    def _is_real_topic(value: str) -> bool:
        normalized = str(value or "").strip()
        compact = "".join(normalized.split())
        if len(normalized) < 2:
            return False
        if re.fullmatch(r"\d+\s*(?:道|个|题)|[一二三四五六七八九十百]+\s*(?:道|个|题)", compact):
            return False
        placeholders = {
            "生成PPT",
            "生成ppt",
            "生成一份PPT",
            "生成一份ppt",
            "一份",
            "一个",
            "一种",
            "一套",
            "一些",
            "PPT",
            "ppt",
            "PPT大纲",
            "ppt大纲",
            "PPT文件",
            "ppt文件",
            "大纲",
            "课件",
            "幻灯片",
            "演示文稿",
            "生成文档",
            "生成练习题",
            "给我一份文档",
            "给我一套资源",
            "给我生成资源",
            "出5道题",
            "出五道题",
            "出题",
            "我当前想学的主题",
            "根据当前阶段",
            "根据当前学习阶段",
            "当前阶段",
            "当前学习阶段",
            "当前主题",
            "学习主题",
            "主题",
        }
        if normalized in placeholders or compact in placeholders:
            return False
        if re.fullmatch(
            r"(?:一份|一个|一种|一套|一些|几个|几道)?(?:PPT|ppt|slides?|PPT大纲|ppt大纲|PPT文件|ppt文件|大纲|课件|幻灯片|演示文稿|文档|资料|资源|学习资源|练习题|习题|题目|视频|代码案例)",
            compact,
        ):
            return False
        return not ResourceGenerationService._looks_like_resource_command(normalized)

    @staticmethod
    def _looks_like_resource_command(value: str) -> bool:
        text = str(value or "").strip()
        if not text:
            return True
        has_resource_word = any(
            word in text.lower()
            for word in (
                "ppt",
                "slides",
                "文档",
                "讲义",
                "资源",
                "练习题",
                "习题",
                "题目",
                "思维导图",
                "导图",
                "视频",
                "代码案例",
            )
        )
        has_action = any(
            word in text
            for word in ("生成", "制作", "创建", "整理", "准备", "设计", "编写", "给我", "帮我", "出")
        )
        cleaned = text
        for word in (
            "请",
            "帮我",
            "给我",
            "生成",
            "制作",
            "创建",
            "整理",
            "准备",
            "设计",
            "编写",
            "一个",
            "一种",
            "一份",
            "一套",
            "一些",
            "几个",
            "几道",
            "文档",
            "讲义",
            "ppt",
            "PPT",
            "slides",
            "幻灯片",
            "演示文稿",
            "课件",
            "练习题",
            "习题",
            "题目",
            "资源",
            "学习资源",
            "资料",
        ):
            cleaned = cleaned.replace(word, "")
        cleaned = re.sub(r"\d+\s*道|[一二三四五六七八九十百]+\s*道", "", cleaned)
        cleaned = cleaned.strip(" ：:，,。.、！？!?")
        return has_resource_word and has_action and len(cleaned) < 2

    async def _build_mindmap(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        display_topic = self._display_topic(params)
        title = f"{display_topic}思维导图"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        generated_mindmap = await self.content_chain.generate_mindmap_asset(
            title=title,
            topic=display_topic,
            snapshot=generation_snapshot,
            sources=sources,
        )
        mermaid = self._render_mindmap_mermaid(generated_mindmap)
        file_name = self._scoped_file_name("mindmap", "mmd", params)
        path = self._write_text(file_name, mermaid)
        return GeneratedAsset(
            assetType="MINDMAP",
            title=generated_mindmap.title,
            summary=generated_mindmap.summary,
            displayMode="INLINE_MERMAID",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_mindmap.title,
            mimeType="text/plain; charset=UTF-8",
            inlineContent=mermaid,
        )

    def _render_mindmap_mermaid(self, generated_mindmap: GeneratedMindMap) -> str:
        lines = ["mindmap", f'  root["{self._escape_mermaid_label(generated_mindmap.root)}"]']
        node_index = 0

        def append_nodes(nodes: list[Any], depth: int) -> None:
            nonlocal node_index
            indent = "  " * depth
            for node in nodes:
                node_index += 1
                lines.append(
                    f'{indent}node_{node_index}["{self._escape_mermaid_label(str(node.name))}"]'
                )
                append_nodes(node.children, depth + 1)

        append_nodes(generated_mindmap.children, 2)
        return "\n".join(lines) + "\n"

    @staticmethod
    def _escape_mermaid_label(label: str) -> str:
        return (
            str(label)
            .replace("\\", "\\\\")
            .replace('"', '\\"')
            .replace("\r", " ")
            .replace("\n", " ")
            .strip()
        )

    async def _build_code(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        display_topic = self._display_topic(params)
        title = f"{display_topic}代码案例"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        generated_code = await self.content_chain.generate_code_asset(
            title=title,
            topic=display_topic,
            snapshot=generation_snapshot,
            sources=sources,
        )
        code_suffix = self._code_file_suffix(generated_code.language)
        file_name = self._scoped_file_name("code_case", code_suffix, params)
        path = self._write_text(file_name, generated_code.code)
        return GeneratedAsset(
            assetType="CODE",
            title=generated_code.title,
            summary=generated_code.summary,
            displayMode="INLINE_CODE",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_code.title,
            mimeType="text/plain; charset=UTF-8",
            inlineContent=generated_code.code,
            language=generated_code.language,
            explanation=generated_code.explanation,
        )

    def _safe_task_id(self, task_id: str) -> str:
        return "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in task_id)

    def _normalize_duration_seconds(self, value: Any) -> int:
        try:
            duration = int(value)
        except (TypeError, ValueError):
            duration = 60
        return max(15, min(180, duration))

    def _build_video_thumbnail_svg(self, *, title: str, topic: str, style: str) -> str:
        style_label = {
            "talking_head": "数字人讲解",
            "animation": "动画演示",
            "hybrid": "混合讲解",
        }.get(style, "教学视频")
        safe_title = self._escape_svg_text(title)
        safe_topic = self._escape_svg_text(topic)
        safe_style = self._escape_svg_text(style_label)
        return (
            '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">'
            '<defs><linearGradient id="bg" x1="0" y1="0" x2="1" y2="1">'
            '<stop offset="0%" stop-color="#0f172a"/><stop offset="100%" stop-color="#2563eb"/>'
            '</linearGradient></defs>'
            '<rect width="1280" height="720" fill="url(#bg)"/>'
            '<rect x="64" y="64" width="1152" height="592" rx="32" fill="rgba(15,23,42,0.28)" stroke="#93c5fd" stroke-width="2"/>'
            '<text x="100" y="160" fill="#cbd5e1" font-size="30" font-family="Segoe UI, Arial, sans-serif">AI Teaching Video</text>'
            f'<text x="100" y="268" fill="#ffffff" font-size="58" font-family="Segoe UI, Arial, sans-serif">{safe_title}</text>'
            f'<text x="100" y="352" fill="#dbeafe" font-size="34" font-family="Segoe UI, Arial, sans-serif">知识点: {safe_topic}</text>'
            f'<text x="100" y="412" fill="#bfdbfe" font-size="30" font-family="Segoe UI, Arial, sans-serif">风格: {safe_style}</text>'
            '<circle cx="1070" cy="360" r="84" fill="#ffffff" fill-opacity="0.9"/>'
            '<polygon points="1042,314 1042,406 1118,360" fill="#2563eb"/>'
            "</svg>"
        )

    def _escape_svg_text(self, value: str) -> str:
        return (
            value.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
        )

    def _write_text(self, file_name: str, content: str) -> Path:
        self.sandbox_root.mkdir(parents=True, exist_ok=True)
        path = self.sandbox_root / file_name
        path.write_text(content, encoding="utf-8")
        return path

    @staticmethod
    def _code_file_suffix(language: str | None) -> str:
        normalized = (language or "").strip().lower()
        suffix_by_language = {
            "python": "py",
            "py": "py",
            "javascript": "js",
            "js": "js",
            "typescript": "ts",
            "ts": "ts",
            "java": "java",
            "rust": "rs",
            "rs": "rs",
            "go": "go",
            "c": "c",
            "cpp": "cpp",
            "c++": "cpp",
            "sql": "sql",
            "html": "html",
            "css": "css",
            "shell": "sh",
            "bash": "sh",
        }
        return suffix_by_language.get(normalized, "txt")

    def _scoped_file_name(self, prefix: str, suffix: str, params: dict[str, Any]) -> str:
        task_id = str(params.get("taskId") or "shared")
        safe_task_id = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in task_id)
        return f"{prefix}_{safe_task_id}.{suffix}"
